import os
import json
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
from typing import List, Dict

class KnowledgeBase:
    def __init__(self, base_dir: str = "knowledge_base"):
        self.base_dir = base_dir
        self.knowledge = {}
        self._init_knowledge_base()
        
    def _init_knowledge_base(self):
        """初始化知识库目录"""
        if not os.path.exists(self.base_dir):
            os.makedirs(self.base_dir)
            
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """从PDF文件提取文本"""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
        return text
        
    def _extract_text_from_docx(self, file_path: str) -> str:
        """从DOCX文件提取文本"""
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        
    def _extract_text_from_pptx(self, file_path: str) -> str:
        """从PPTX文件提取文本"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
        
    def _extract_text_from_excel(self, file_path: str) -> str:
        """从Excel文件提取文本"""
        df = pd.read_excel(file_path)
        return df.to_string()
        
    def add_file(self, file_path: str) -> bool:
        """添加文件到知识库"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            target_path = os.path.join(self.base_dir, os.path.basename(file_path))
            
            # 复制文件到知识库目录
            with open(file_path, 'rb') as src, open(target_path, 'wb') as dst:
                dst.write(src.read())
            
            # 提取文本内容
            content = ""
            if ext == '.pdf':
                content = self._extract_text_from_pdf(target_path)
            elif ext == '.docx':
                content = self._extract_text_from_docx(target_path)
            elif ext == '.pptx':
                content = self._extract_text_from_pptx(target_path)
            elif ext == '.xlsx' or ext == '.xls':
                content = self._extract_text_from_excel(target_path)
            elif ext == '.txt':
                with open(target_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    
            self.knowledge[os.path.basename(file_path)] = content
            return True
            
        except Exception as e:
            print(f"添加文件失败: {str(e)}")
            return False
            
    def get_relevant_knowledge(self, query: str) -> str:
        """获取与查询相关的知识"""
        # 这里可以实现更复杂的相关性搜索算法
        relevant_content = []
        for filename, content in self.knowledge.items():
            if query.lower() in content.lower():
                relevant_content.append(f"从 {filename} 中找到相关内容：\n{content[:500]}...")
        
        return "\n\n".join(relevant_content) if relevant_content else ""